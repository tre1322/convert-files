# converters/pptx.py
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
import tempfile

def pdf_to_pptx(in_path, out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)  # 16:9
    pages = convert_from_path(str(in_path), dpi=200)
    for img in pages:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            img.save(tmp.name, "PNG")
            slide.shapes.add_picture(tmp.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out_path)
